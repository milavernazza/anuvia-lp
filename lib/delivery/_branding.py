"""Anuvia delivery branding — shared visual identity for client deliverables.

Single source of truth for HTML/PDF + PPTX rendering across all delivery
agents (finops_audit, ai_readiness, devops_maturity, ...). Tied to the LP
brand tokens defined in ``templates/_base.html``.

Public API:

  * ``render_deliverable_html(...)`` — full HTML document with dark-banner
    cover page, section dividers, callout boxes, stat cards, running
    headers/footers, table styling, blockquotes, code blocks, and inline
    SVG charts.
  * ``md_to_html_rich(md)`` — markdown -> HTML using ``markdown`` package
    with the ``tables``, ``fenced_code``, ``nl2br`` extensions, then a
    multi-pass post-processing step to bolt Anuvia classnames onto tables,
    blockquotes, detect FinOps framework metadata blocks, render callouts
    for "Premissas e Limitações" / "Validation criteria", and inject
    inline SVG bar/donut charts when a structured "chart:" marker is
    present.
  * ``svg_bar_chart(items, ...)`` — horizontal range bar chart helper.
  * ``svg_donut_chart(items, ...)`` — single-ring donut helper.
  * ``generate_pptx_deck(...)`` — build a PPTX deck in memory with the
    new dark/light cover, full-bleed section dividers, and content slides
    with a slim top band; returns the binary bytes.

Design rules — DO NOT INVENT colors/fonts. Strict monochrome palette pulled
directly from ``templates/_base.html`` (#1a1a1a, #fafaf9, #e7e5e4, #78716c,
white). Functional accents (red/blue/green) reserved exclusively for the
warning/info/success callout variants — never used decoratively.
"""

from __future__ import annotations

import io
import math
import re
from datetime import datetime, timezone
from typing import Any, Dict, List, Optional, Tuple


# ---------------------------------------------------------------------------
# Brand tokens — single source of truth (mirror of templates/_base.html)
# ---------------------------------------------------------------------------

BG_PAPER = "#fafaf9"          # warm off-white (page bg, card bg)
INK = "#1a1a1a"               # near-black (primary text, headings, dark banners)
INK_MUTED = "#475569"         # secondary text on dark bg
STONE_500 = "#78716c"         # muted text (eyebrow)
STONE_400 = "#94a3b8"         # cool muted (subtitles on dark bg)
STONE_200 = "#e7e5e4"         # borders / rules
STONE_100 = "#f5f5f4"         # alt panel bg
WHITE = "#ffffff"

# Functional accents — used ONLY for callout box borders. Never decorative.
ACCENT_WARN = "#b91c1c"       # warning callout
ACCENT_INFO = "#1d4ed8"       # info callout
ACCENT_OK = "#15803d"         # success callout

FONT_BODY = "'Inter', system-ui, -apple-system, 'Helvetica Neue', sans-serif"
FONT_HEAD = "'Playfair Display', Georgia, serif"
FONT_MONO = "'JetBrains Mono', 'SF Mono', ui-monospace, monospace"

GOOGLE_FONTS_LINK = (
    "https://fonts.googleapis.com/css2?"
    "family=Inter:wght@400;500;600;700"
    "&family=Playfair+Display:wght@400;500;600;700"
    "&family=JetBrains+Mono:wght@400;500&display=swap"
)


def _now() -> datetime:
    return datetime.now(timezone.utc)


def _today_pt() -> str:
    return _now().strftime("%d/%m/%Y")


# ---------------------------------------------------------------------------
# Markdown -> HTML rich
# ---------------------------------------------------------------------------


def md_to_html_rich(md: str) -> str:
    """Convert markdown to HTML with Anuvia styling.

    Uses the ``markdown`` package with ``tables``, ``fenced_code``, ``nl2br``
    so GFM-style tables actually render. Falls back to a tiny hand-rolled
    converter if the package is unavailable (keeps the delivery pipeline
    alive on minimal images).
    """
    if not md:
        return ""

    try:
        import markdown  # type: ignore
    except Exception:  # noqa: BLE001
        return _md_fallback(md)

    try:
        html = markdown.markdown(
            md,
            extensions=["tables", "fenced_code", "nl2br", "sane_lists"],
            output_format="html5",
        )
    except Exception:  # noqa: BLE001
        # Re-try without sane_lists, which some older markdown versions
        # don't ship.
        try:
            html = markdown.markdown(
                md,
                extensions=["tables", "fenced_code", "nl2br"],
                output_format="html5",
            )
        except Exception:  # noqa: BLE001
            return _md_fallback(md)

    # Bolt Anuvia classnames onto the generated tags so the CSS hooks line up.
    html = re.sub(r"<table>", '<table class="anuvia-table">', html)
    html = re.sub(r"<blockquote>", '<blockquote class="anuvia-bq">', html)

    # FinOps framework polish — promote special H2 sections to callout boxes
    # and wrap "## Premissas e Limitações" / similar in the styled box.
    html = _wrap_callouts(html)

    # Detect inline chart markers like ``[chart:bar] ...`` and inject SVG.
    html = _inject_inline_charts(html)

    return html


# Callouts and chart-marker post-processors (run after `markdown` package
# converts md -> html). These look for predictable patterns produced by the
# Claude prompts and don't rely on the markdown extension API, so the
# pipeline stays portable.

_CALLOUT_TITLES = {
    "premissas e limitações": "info",
    "premissas e limitacoes": "info",
    "premissas": "info",
    "validation criteria": "info",
    "critérios de validação": "info",
    "criterios de validacao": "info",
    "confiança baixa": "warning",
    "confianca baixa": "warning",
    "atenção": "warning",
    "atencao": "warning",
    "próximos passos": "success",
    "proximos passos": "success",
    "próxima rodada — o que pedir": "warning",
    "proxima rodada - o que pedir": "warning",
    "próxima rodada": "warning",
}


def _wrap_callouts(html: str) -> str:
    """Wrap selected H2 sections in callout boxes.

    Heuristic: for each ``<h2>...</h2>`` whose text matches a known key,
    rewrap the H2 plus the following sibling block(s) up to the next H2
    inside a ``<aside class="callout callout-...">`` element. We do this
    with a regex pass (markdown output is predictable, single-line tags).
    """
    if "<h2" not in html:
        return html

    # Split by H2 boundaries to make wrapping tractable.
    parts = re.split(r"(<h2[^>]*>.*?</h2>)", html, flags=re.DOTALL)
    out: List[str] = []
    i = 0
    while i < len(parts):
        chunk = parts[i]
        if chunk.startswith("<h2") and i + 1 < len(parts):
            # The next part is the body until the next H2 marker.
            body = parts[i + 1] if i + 1 < len(parts) else ""
            # Find h2 text.
            m = re.search(r"<h2[^>]*>(.*?)</h2>", chunk, flags=re.DOTALL)
            title_text = (m.group(1) if m else "").strip().lower()
            # Strip any nested tags.
            title_text = re.sub(r"<[^>]+>", "", title_text).strip()
            kind = _CALLOUT_TITLES.get(title_text)
            if kind:
                label = "PREMISSA" if kind == "info" else (
                    "ATENÇÃO" if kind == "warning" else "PRÓXIMOS"
                )
                # Strip the H2 itself; we'll render a custom label instead.
                inner = body
                out.append(
                    f'<aside class="callout callout-{kind}">'
                    f'<div class="callout-label">{label}</div>'
                    f'<div class="callout-title">{_html_escape(_title_case(title_text))}</div>'
                    f'<div class="callout-body">{inner}</div>'
                    f'</aside>'
                )
                i += 2
                continue
            out.append(chunk)
            out.append(body)
            i += 2
            continue
        out.append(chunk)
        i += 1
    return "".join(out)


def _title_case(s: str) -> str:
    # Keep accents — just upper-first each significant word.
    return s.title() if s else s


# Inline chart marker: a paragraph that contains exactly ``[chart:bar:JSON]``
# or ``[chart:donut:JSON]`` is replaced with the SVG. JSON is a list of
# ``[label, low, high]`` (bar) or ``[label, value]`` (donut). This lets the
# Claude prompts request a chart without coupling to matplotlib.
_CHART_RE = re.compile(
    r"<p>\s*\[chart:(bar|donut):(.+?)\]\s*</p>",
    re.DOTALL,
)


def _inject_inline_charts(html: str) -> str:
    def _sub(m: re.Match) -> str:
        kind = m.group(1)
        raw = m.group(2)
        try:
            import json as _json
            data = _json.loads(raw)
        except Exception:  # noqa: BLE001
            return m.group(0)
        if not isinstance(data, list):
            return m.group(0)
        if kind == "bar":
            items: List[Tuple[str, float, float]] = []
            for row in data:
                if isinstance(row, list) and len(row) >= 3:
                    items.append((str(row[0]), float(row[1]), float(row[2])))
                elif isinstance(row, list) and len(row) == 2:
                    items.append((str(row[0]), float(row[1]), float(row[1])))
            return svg_bar_chart(items)
        if kind == "donut":
            ditems: List[Tuple[str, float]] = []
            for row in data:
                if isinstance(row, list) and len(row) >= 2:
                    ditems.append((str(row[0]), float(row[1])))
            return svg_donut_chart(ditems)
        return m.group(0)

    return _CHART_RE.sub(_sub, html)


def _md_fallback(md: str) -> str:
    """Last-resort tiny converter — preserves headings, lists, bold, code.

    Used only if the ``markdown`` package import fails. Matches the older
    hand-rolled converter so existing tests pass.
    """
    lines: List[str] = []
    in_list = False
    for raw in md.splitlines():
        line = raw.rstrip()
        if not line.strip():
            if in_list:
                lines.append("</ul>")
                in_list = False
            lines.append("")
            continue
        m = re.match(r"^(#{1,4})\s+(.*)$", line.strip())
        if m:
            if in_list:
                lines.append("</ul>")
                in_list = False
            level = len(m.group(1))
            tag = f"h{min(level + 1, 6)}"
            lines.append(f"<{tag}>{m.group(2)}</{tag}>")
            continue
        m = re.match(r"^\s*[-*]\s+(.*)$", line)
        if m:
            if not in_list:
                lines.append("<ul>")
                in_list = True
            content = m.group(1)
            content = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", content)
            content = re.sub(r"`([^`]+)`", r"<code>\1</code>", content)
            lines.append(f"<li>{content}</li>")
            continue
        if in_list:
            lines.append("</ul>")
            in_list = False
        content = re.sub(r"\*\*(.+?)\*\*", r"<strong>\1</strong>", line)
        content = re.sub(r"`([^`]+)`", r"<code>\1</code>", content)
        lines.append(f"<p>{content}</p>")
    if in_list:
        lines.append("</ul>")
    return "\n".join(lines)


# ---------------------------------------------------------------------------
# HTML deliverable rendering
# ---------------------------------------------------------------------------


def _html_escape(s: str) -> str:
    if s is None:
        return ""
    return (
        str(s)
        .replace("&", "&amp;")
        .replace("<", "&lt;")
        .replace(">", "&gt;")
        .replace('"', "&quot;")
    )


# ---------------------------------------------------------------------------
# Stat card + section divider builders
# ---------------------------------------------------------------------------


def stat_card_html(label: str, value: str, meta: Optional[str] = None) -> str:
    """Single stat card — eyebrow label + Playfair value + optional meta."""
    meta_html = (
        f'<div class="stat-meta">{_html_escape(meta)}</div>' if meta else ""
    )
    return (
        '<div class="stat-card">'
        f'<div class="stat-label">{_html_escape(label)}</div>'
        f'<div class="stat-value">{_html_escape(value)}</div>'
        f'{meta_html}'
        '</div>'
    )


def stat_grid_html(stats: List[Tuple[str, str, Optional[str]]]) -> str:
    """Stat grid wrapper — list of (label, value, meta) tuples."""
    if not stats:
        return ""
    cards = "".join(stat_card_html(*s) for s in stats)
    return f'<div class="stat-grid">{cards}</div>'


def section_divider_html(section_num: int, title: str, subtitle: str = "") -> str:
    """Full-page section divider — slate bg, white centered Playfair title."""
    sub_html = (
        f'<div class="section-divider-subtitle">{_html_escape(subtitle)}</div>'
        if subtitle else ""
    )
    return (
        '<section class="section-divider">'
        f'<div class="section-divider-eyebrow">SEÇÃO {section_num:02d}</div>'
        f'<div class="section-divider-title">{_html_escape(title)}</div>'
        f'{sub_html}'
        '</section>'
    )


# ---------------------------------------------------------------------------
# Inline SVG charts (no matplotlib dependency)
# ---------------------------------------------------------------------------


def svg_bar_chart(
    items: List[Tuple[str, float, float]],
    *,
    width: int = 620,
    height: int = 260,
    currency: str = "R$",
) -> str:
    """Horizontal range bar chart.

    Args:
        items: list of ``(label, value_low, value_high)`` — values in
            same currency unit. ``value_low`` may equal ``value_high`` for
            single-point series.
        width / height: SVG canvas size in px.
        currency: prefix for axis labels.

    Strict monochrome: dark INK bars (low), STONE_500 extension (high),
    STONE_200 grid. Inter 11px labels.
    """
    if not items:
        return ""

    n = len(items)
    margin_left = 160
    margin_right = 90
    margin_top = 18
    margin_bottom = 22
    plot_w = max(width - margin_left - margin_right, 100)
    plot_h = max(height - margin_top - margin_bottom, 40)
    row_h = plot_h / n

    max_val = max((v[2] for v in items), default=1.0) or 1.0
    # Round up max to a "nice" gridline (1, 2, 5 × 10^n).
    grid_max = _nice_ceil(max_val)
    gridlines = 4

    def _x(v: float) -> float:
        return margin_left + (v / grid_max) * plot_w

    def _fmt(v: float) -> str:
        if grid_max >= 1_000_000:
            return f"{currency} {v / 1_000_000:.1f}M"
        if grid_max >= 1_000:
            return f"{currency} {v / 1_000:.0f}k"
        return f"{currency} {v:.0f}"

    svg: List[str] = []
    svg.append(
        f'<svg class="anuvia-chart" viewBox="0 0 {width} {height}" '
        f'xmlns="http://www.w3.org/2000/svg" role="img" '
        f'aria-label="Anuvia bar chart">'
    )
    # Gridlines + axis ticks.
    for i in range(gridlines + 1):
        x = margin_left + (plot_w * i / gridlines)
        v = grid_max * i / gridlines
        svg.append(
            f'<line x1="{x:.1f}" y1="{margin_top}" x2="{x:.1f}" '
            f'y2="{margin_top + plot_h}" stroke="{STONE_200}" stroke-width="1"/>'
        )
        svg.append(
            f'<text x="{x:.1f}" y="{margin_top + plot_h + 14}" '
            f'font-family="Inter, sans-serif" font-size="9" fill="{STONE_500}" '
            f'text-anchor="middle">{_fmt(v)}</text>'
        )
    # Rows.
    for idx, (label, lo, hi) in enumerate(items):
        y0 = margin_top + row_h * idx + row_h * 0.18
        bar_h = row_h * 0.55
        x_lo = _x(max(lo, 0))
        x_hi = _x(max(hi, 0))
        # Label (left).
        svg.append(
            f'<text x="{margin_left - 8:.1f}" y="{y0 + bar_h * 0.72:.1f}" '
            f'font-family="Inter, sans-serif" font-size="10.5" '
            f'fill="{INK}" text-anchor="end" font-weight="500">'
            f'{_html_escape(label)}</text>'
        )
        if hi > lo:
            # Range from low to high — INK low bar + STONE_500 extension.
            svg.append(
                f'<rect x="{margin_left:.1f}" y="{y0:.1f}" '
                f'width="{x_hi - margin_left:.1f}" height="{bar_h:.1f}" '
                f'fill="{STONE_500}"/>'
            )
            svg.append(
                f'<rect x="{margin_left:.1f}" y="{y0:.1f}" '
                f'width="{x_lo - margin_left:.1f}" height="{bar_h:.1f}" '
                f'fill="{INK}"/>'
            )
            # Value label.
            svg.append(
                f'<text x="{x_hi + 6:.1f}" y="{y0 + bar_h * 0.72:.1f}" '
                f'font-family="Inter, sans-serif" font-size="10" '
                f'fill="{STONE_500}">{_fmt(lo)}–{_fmt(hi)}</text>'
            )
        else:
            svg.append(
                f'<rect x="{margin_left:.1f}" y="{y0:.1f}" '
                f'width="{x_hi - margin_left:.1f}" height="{bar_h:.1f}" '
                f'fill="{INK}"/>'
            )
            svg.append(
                f'<text x="{x_hi + 6:.1f}" y="{y0 + bar_h * 0.72:.1f}" '
                f'font-family="Inter, sans-serif" font-size="10" '
                f'fill="{STONE_500}">{_fmt(hi)}</text>'
            )

    svg.append("</svg>")
    return "".join(svg)


def svg_donut_chart(
    items: List[Tuple[str, float]],
    *,
    width: int = 320,
    height: int = 320,
) -> str:
    """Donut chart for category breakdown.

    Args:
        items: ``[(label, value), ...]``. Values summed for percentage.

    Strict monochrome: 5 gray shades cycled. Center shows total.
    """
    if not items:
        return ""

    total = sum(max(v, 0) for _, v in items) or 1.0
    cx = width / 2
    cy = height / 2
    r_outer = min(width, height) * 0.38
    r_inner = r_outer * 0.62
    # Gray ramp dark→light.
    shades = ["#1a1a1a", "#3d3d3d", "#6b6b6b", "#9c9a98", "#cfcecd"]

    svg: List[str] = []
    svg.append(
        f'<svg class="anuvia-chart" viewBox="0 0 {width} {height}" '
        f'xmlns="http://www.w3.org/2000/svg" role="img" '
        f'aria-label="Anuvia donut chart">'
    )

    start_angle = -math.pi / 2  # start at top
    for i, (label, value) in enumerate(items):
        if value <= 0:
            continue
        frac = value / total
        end_angle = start_angle + frac * 2 * math.pi
        large_arc = 1 if frac > 0.5 else 0
        # Outer arc points.
        x1 = cx + r_outer * math.cos(start_angle)
        y1 = cy + r_outer * math.sin(start_angle)
        x2 = cx + r_outer * math.cos(end_angle)
        y2 = cy + r_outer * math.sin(end_angle)
        # Inner arc points (reverse direction).
        x3 = cx + r_inner * math.cos(end_angle)
        y3 = cy + r_inner * math.sin(end_angle)
        x4 = cx + r_inner * math.cos(start_angle)
        y4 = cy + r_inner * math.sin(start_angle)
        color = shades[i % len(shades)]
        path = (
            f'M {x1:.1f} {y1:.1f} '
            f'A {r_outer:.1f} {r_outer:.1f} 0 {large_arc} 1 {x2:.1f} {y2:.1f} '
            f'L {x3:.1f} {y3:.1f} '
            f'A {r_inner:.1f} {r_inner:.1f} 0 {large_arc} 0 {x4:.1f} {y4:.1f} '
            f'Z'
        )
        svg.append(f'<path d="{path}" fill="{color}"/>')
        start_angle = end_angle

    # Center label.
    svg.append(
        f'<text x="{cx:.1f}" y="{cy - 6:.1f}" '
        f'font-family="Inter, sans-serif" font-size="9" '
        f'fill="{STONE_500}" text-anchor="middle" letter-spacing="0.16em">TOTAL</text>'
    )
    svg.append(
        f'<text x="{cx:.1f}" y="{cy + 14:.1f}" '
        f'font-family="Playfair Display, serif" font-size="18" '
        f'fill="{INK}" text-anchor="middle" font-weight="600">'
        f'{_fmt_thousands(total)}</text>'
    )

    # Legend below.
    legend_y = height - 8 - 14 * len(items)
    # If chart is small, render legend to the right instead.
    if legend_y < r_outer + cy + 4:
        legend_x = cx + r_outer + 14
        for i, (label, value) in enumerate(items):
            color = shades[i % len(shades)]
            ly = (cy - r_outer) + i * 16
            svg.append(
                f'<rect x="{legend_x:.1f}" y="{ly:.1f}" width="9" height="9" '
                f'fill="{color}"/>'
            )
            pct = (value / total) * 100 if total else 0
            svg.append(
                f'<text x="{legend_x + 14:.1f}" y="{ly + 8:.1f}" '
                f'font-family="Inter, sans-serif" font-size="10" '
                f'fill="{INK}">{_html_escape(label)} '
                f'<tspan fill="{STONE_500}">({pct:.0f}%)</tspan></text>'
            )
    svg.append("</svg>")
    return "".join(svg)


def _nice_ceil(v: float) -> float:
    """Round up to a 'nice' tick boundary (1, 2, 5 × 10^n)."""
    if v <= 0:
        return 1.0
    exp = math.floor(math.log10(v))
    base = 10 ** exp
    for mult in (1, 2, 2.5, 5, 10):
        candidate = mult * base
        if candidate >= v:
            return candidate
    return 10 * base


def _fmt_thousands(v: float) -> str:
    if v >= 1_000_000:
        return f"R$ {v / 1_000_000:.2f}M"
    if v >= 1_000:
        return f"R$ {v / 1_000:.0f}k"
    return f"R$ {v:.0f}"


def _derive_cover_stats(
    meta: Dict[str, Any],
) -> List[Tuple[str, str, Optional[str]]]:
    """Derive stat cards from engagement meta when caller didn't pass any.

    Heuristic: look for canonical Anuvia FinOps meta keys and pull the
    headline numbers out. Returns up to 4 cards.
    """
    if not meta:
        return []
    stats: List[Tuple[str, str, Optional[str]]] = []
    if meta.get("Baseline mensal"):
        stats.append(("Baseline mensal", str(meta["Baseline mensal"]), None))
    if meta.get("Economia identificada"):
        stats.append(
            ("Economia anualizada", str(meta["Economia identificada"]), "faixa estimada")
        )
    if meta.get("Payback"):
        stats.append(("Payback", str(meta["Payback"]), None))
    if meta.get("Período"):
        stats.append(("Período", str(meta["Período"]), None))
    return stats[:4]


def _cover_meta_block(meta: Dict[str, Any]) -> str:
    """Render the engagement-meta card on the cover page."""
    if not meta:
        return ""
    order = (
        "Cliente",
        "Período",
        "Baseline mensal",
        "Economia identificada",
        "Payback",
        "Analista responsável",
        "Analista",
    )
    rows: List[str] = []
    seen: set = set()
    for key in order:
        if key in seen or key not in meta:
            continue
        seen.add(key)
        value = meta.get(key)
        if value in (None, ""):
            continue
        rows.append(
            f'<div class="meta-row">'
            f'<div class="meta-label">{_html_escape(key)}</div>'
            f'<div class="meta-value">{_html_escape(value)}</div>'
            f"</div>"
        )
    # Any extra keys not in the canonical order — keep them, but at the end.
    for key, value in meta.items():
        if key in seen or value in (None, ""):
            continue
        seen.add(key)
        rows.append(
            f'<div class="meta-row">'
            f'<div class="meta-label">{_html_escape(key)}</div>'
            f'<div class="meta-value">{_html_escape(value)}</div>'
            f"</div>"
        )
    if not rows:
        return ""
    return f'<div class="meta-card">{"".join(rows)}</div>'


def _stylesheet(practice_label: str) -> str:
    """The full Anuvia stylesheet. Inlined so Gotenberg can render offline."""
    label_esc = _html_escape(practice_label)
    return f"""
@import url('{GOOGLE_FONTS_LINK}');

@page {{
  size: A4;
  margin: 22mm 18mm 22mm 18mm;
  @top-left {{
    content: "ANUVIA · {label_esc}";
    font-family: {FONT_BODY};
    font-size: 9px;
    letter-spacing: 0.18em;
    color: {STONE_500};
    font-weight: 600;
    text-transform: uppercase;
  }}
  @top-right {{
    content: "";
    border-bottom: 1px solid {STONE_200};
  }}
  @bottom-left {{
    content: "Anuvia Cloud & AI Consulting · Mila Vernazza";
    font-family: {FONT_BODY};
    font-size: 9px;
    color: {STONE_500};
  }}
  @bottom-right {{
    content: counter(page) " / " counter(pages);
    font-family: {FONT_BODY};
    font-size: 9px;
    color: {STONE_500};
  }}
}}

@page :first {{
  @top-left {{ content: ""; }}
  @top-right {{ content: ""; }}
  @bottom-left {{ content: ""; }}
  @bottom-right {{ content: ""; }}
}}

* {{ box-sizing: border-box; }}

html, body {{
  margin: 0;
  padding: 0;
  background: {WHITE};
  color: {INK};
  font-family: {FONT_BODY};
  font-size: 11px;
  line-height: 1.65;
  -webkit-font-smoothing: antialiased;
}}

.eyebrow {{
  font-family: {FONT_BODY};
  font-size: 11px;
  letter-spacing: 0.18em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_500};
}}

/* -------- Cover page (dark banner + white body) -------- */

.cover-page {{
  page-break-after: always;
  min-height: 100vh;
  padding: 0;
  margin: -22mm -18mm 0 -18mm;  /* bleed full A4 width past the @page margin */
  background: {WHITE};
}}

.cover-banner {{
  background: {INK};
  color: {WHITE};
  padding: 28mm 22mm 24mm 22mm;
  position: relative;
  overflow: hidden;
}}

.cover-banner::after {{
  /* Subtle diagonal accent — three thin monochrome lines top-right corner. */
  content: "";
  position: absolute;
  top: -40px;
  right: -40px;
  width: 200px;
  height: 200px;
  background:
    linear-gradient(45deg, transparent 49%, rgba(255,255,255,0.06) 49%, rgba(255,255,255,0.06) 51%, transparent 51%) center / 100% 100%,
    linear-gradient(45deg, transparent 64%, rgba(255,255,255,0.04) 64%, rgba(255,255,255,0.04) 66%, transparent 66%) center / 100% 100%,
    linear-gradient(45deg, transparent 79%, rgba(255,255,255,0.03) 79%, rgba(255,255,255,0.03) 81%, transparent 81%) center / 100% 100%;
  transform: rotate(0deg);
  pointer-events: none;
}}

.cover-banner-top {{
  display: flex;
  justify-content: space-between;
  align-items: baseline;
  margin-bottom: 22mm;
  position: relative;
  z-index: 1;
}}

.cover-wordmark {{
  font-family: {FONT_HEAD};
  font-size: 18pt;
  font-weight: 600;
  letter-spacing: 0.04em;
  color: {WHITE};
}}

.cover-banner-eyebrow {{
  font-family: {FONT_BODY};
  font-size: 9pt;
  letter-spacing: 0.22em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_400};
}}

.cover-title {{
  font-family: {FONT_HEAD};
  font-weight: 600;
  font-size: 48pt;
  line-height: 1.02;
  letter-spacing: -0.02em;
  color: {WHITE};
  margin: 0 0 14px;
  position: relative;
  z-index: 1;
}}

.cover-subtitle {{
  font-family: {FONT_BODY};
  font-size: 14pt;
  color: {STONE_400};
  margin: 0 0 4px;
  font-weight: 400;
  position: relative;
  z-index: 1;
}}

.cover-engagement-id {{
  font-family: {FONT_MONO};
  font-size: 10pt;
  color: {STONE_400};
  letter-spacing: 0.06em;
  position: relative;
  z-index: 1;
}}

.cover-body {{
  padding: 18mm 22mm 14mm 22mm;
  display: grid;
  grid-template-columns: 1fr 1fr;
  gap: 24px;
}}

.cover-body-eyebrow {{
  font-family: {FONT_BODY};
  font-size: 9pt;
  letter-spacing: 0.22em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_500};
  margin: 0 0 10px;
}}

.meta-card {{
  background: {WHITE};
  border: 1px solid {STONE_200};
  padding: 18px 22px;
  display: block;
}}

.meta-row {{
  display: flex;
  align-items: baseline;
  gap: 16px;
  padding: 8px 0;
  border-bottom: 1px solid {STONE_200};
}}

.meta-row:last-child {{
  border-bottom: none;
}}

.meta-label {{
  flex: 0 0 42%;
  font-size: 9.5px;
  letter-spacing: 0.18em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_500};
}}

.meta-value {{
  flex: 1 1 auto;
  font-size: 11.5px;
  color: {INK};
  font-weight: 500;
}}

/* Stat grid on the cover (right column). */

.cover-body .stat-grid {{
  grid-template-columns: repeat(2, 1fr);
}}

.stat-grid {{
  display: grid;
  grid-template-columns: repeat(2, 1fr);
  gap: 12px;
}}

.stat-card {{
  border: 1px solid {STONE_200};
  background: {WHITE};
  padding: 14px 16px;
}}

.stat-label {{
  font-family: {FONT_BODY};
  font-size: 9px;
  letter-spacing: 0.18em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_500};
  margin-bottom: 6px;
}}

.stat-value {{
  font-family: {FONT_HEAD};
  font-size: 22pt;
  font-weight: 600;
  color: {INK};
  letter-spacing: -0.01em;
  line-height: 1.05;
  margin-top: 4px;
}}

.stat-meta {{
  font-family: {FONT_BODY};
  font-size: 10px;
  color: {STONE_500};
  margin-top: 6px;
}}

.cover-bottom {{
  margin: 0 22mm;
  padding: 14px 0 18px 0;
  border-top: 1px solid {STONE_200};
  display: flex;
  justify-content: space-between;
  align-items: center;
  font-size: 9.5px;
  color: {STONE_500};
  letter-spacing: 0.04em;
}}

/* -------- Section divider (full-page break) -------- */

.section-divider {{
  page-break-before: always;
  page-break-after: always;
  margin: -22mm -18mm;  /* bleed full A4 past page margin */
  min-height: 100vh;
  background: {INK};
  color: {WHITE};
  display: flex;
  flex-direction: column;
  align-items: center;
  justify-content: center;
  text-align: center;
  padding: 30mm 22mm;
}}

.section-divider-eyebrow {{
  font-family: {FONT_BODY};
  font-size: 9pt;
  letter-spacing: 0.32em;
  text-transform: uppercase;
  font-weight: 600;
  color: {STONE_400};
  margin-bottom: 22px;
}}

.section-divider-title {{
  font-family: {FONT_HEAD};
  font-size: 42pt;
  font-weight: 600;
  color: {WHITE};
  letter-spacing: -0.015em;
  line-height: 1.05;
  max-width: 28ch;
}}

.section-divider-subtitle {{
  font-family: {FONT_BODY};
  font-size: 13pt;
  color: {STONE_400};
  margin-top: 14px;
  max-width: 40ch;
  font-weight: 400;
}}

/* -------- Callout boxes -------- */

.callout {{
  margin: 18px 0;
  padding: 14px 18px;
  border-left: 4px solid {INK};
  background: {BG_PAPER};
  border-radius: 0 4px 4px 0;
  page-break-inside: avoid;
}}

.callout.callout-warning {{
  border-left-color: {ACCENT_WARN};
}}

.callout.callout-info {{
  border-left-color: {ACCENT_INFO};
}}

.callout.callout-success {{
  border-left-color: {ACCENT_OK};
}}

.callout-label {{
  font-family: {FONT_BODY};
  font-size: 9px;
  letter-spacing: 0.18em;
  text-transform: uppercase;
  font-weight: 700;
  color: {STONE_500};
  margin-bottom: 4px;
}}

.callout-title {{
  font-family: {FONT_HEAD};
  font-size: 14pt;
  font-weight: 600;
  color: {INK};
  margin-bottom: 8px;
  letter-spacing: -0.01em;
}}

.callout-body {{
  font-size: 11px;
}}

.callout-body p:first-child {{ margin-top: 0; }}
.callout-body p:last-child {{ margin-bottom: 0; }}

/* -------- Inline SVG charts -------- */

.anuvia-chart {{
  display: block;
  width: 100%;
  max-width: 620px;
  height: auto;
  margin: 16px 0;
}}

/* -------- Body -------- */

.body-wrap {{
  padding: 0;
}}

.body-eyebrow {{
  margin: 0 0 4px;
}}

.body-title {{
  font-family: {FONT_HEAD};
  font-weight: 600;
  font-size: 28pt;
  letter-spacing: -0.015em;
  color: {INK};
  margin: 0 0 8px;
  line-height: 1.1;
}}

.body-subtitle {{
  font-family: {FONT_BODY};
  font-size: 11pt;
  color: {INK_MUTED};
  margin: 0 0 24px;
  font-weight: 400;
}}

.body-rule {{
  height: 1px;
  background: {STONE_200};
  margin: 0 0 24px;
}}

.body-content h1 {{
  font-family: {FONT_HEAD};
  font-size: 28px;
  font-weight: 600;
  color: {INK};
  margin: 32px 0 8px;
  letter-spacing: -0.015em;
}}

.body-content h2 {{
  font-family: {FONT_HEAD};
  font-size: 18px;
  font-weight: 600;
  color: {INK};
  margin: 32px 0 12px;
  padding-bottom: 6px;
  border-bottom: 1px solid {STONE_200};
  letter-spacing: -0.01em;
}}

.body-content h3 {{
  font-family: {FONT_BODY};
  font-size: 11px;
  font-weight: 600;
  color: {INK};
  text-transform: uppercase;
  letter-spacing: 0.04em;
  margin: 20px 0 8px;
}}

.body-content h4 {{
  font-family: {FONT_BODY};
  font-size: 12px;
  font-weight: 600;
  color: {INK_MUTED};
  margin: 14px 0 6px;
}}

.body-content p {{
  margin: 8px 0 10px;
}}

.body-content ul, .body-content ol {{
  margin: 6px 0 12px;
  padding-left: 22px;
}}

.body-content li {{
  margin: 4px 0;
}}

.body-content li > ul, .body-content li > ol {{
  margin: 4px 0 4px;
}}

.body-content a {{
  color: {INK};
  text-decoration: underline;
  text-decoration-color: {STONE_200};
}}

.body-content strong {{
  font-weight: 600;
  color: {INK};
}}

/* -------- Anuvia table -------- */

.anuvia-table {{
  width: 100%;
  border-collapse: collapse;
  margin: 16px 0;
  font-size: 10.5px;
  page-break-inside: auto;
}}

.anuvia-table thead {{
  display: table-header-group;
}}

.anuvia-table thead th {{
  background: {INK};
  color: {BG_PAPER};
  text-align: left;
  padding: 10px 12px;
  font-weight: 600;
  letter-spacing: 0.02em;
  font-size: 10.5px;
}}

.anuvia-table tbody td {{
  padding: 8px 12px;
  border-bottom: 1px solid {STONE_200};
  vertical-align: top;
}}

.anuvia-table tbody tr:nth-child(even) {{
  background: {BG_PAPER};
}}

.anuvia-table tbody tr {{
  page-break-inside: avoid;
}}

/* -------- Blockquote -------- */

.anuvia-bq {{
  margin: 16px 0;
  padding: 12px 16px;
  border-left: 3px solid {INK};
  background: {BG_PAPER};
  color: {INK_MUTED};
  font-style: italic;
}}

.anuvia-bq p {{
  margin: 4px 0;
}}

/* -------- Code -------- */

code {{
  font-family: {FONT_MONO};
  font-size: 10px;
  background: {BG_PAPER};
  border: 1px solid {STONE_200};
  padding: 1px 5px;
  border-radius: 3px;
}}

pre {{
  background: {INK};
  color: {BG_PAPER};
  padding: 14px 18px;
  border-radius: 6px;
  font-family: {FONT_MONO};
  font-size: 10px;
  line-height: 1.6;
  overflow-x: auto;
  margin: 14px 0;
}}

pre code {{
  background: transparent;
  border: none;
  padding: 0;
  color: inherit;
  font-size: inherit;
}}

hr {{
  border: none;
  border-top: 1px solid {STONE_200};
  margin: 24px 0;
}}

/* Avoid orphan headings */
.body-content h1, .body-content h2, .body-content h3 {{
  page-break-after: avoid;
}}
"""


def render_deliverable_html(
    *,
    practice_label: str,
    title: str,
    subtitle: str,
    body_md: str,
    engagement_meta: Optional[dict] = None,
    show_cover: bool = True,
    cover_stats: Optional[List[Tuple[str, str, Optional[str]]]] = None,
    engagement_id: Optional[str] = None,
) -> str:
    """Return a full HTML document ready for Gotenberg PDF conversion.

    Args:
        practice_label: e.g. ``"FINOPS AUDIT"``. Rendered in eyebrow + footer.
        title: Document title (Playfair big).
        subtitle: Smaller line under the title.
        body_md: Raw markdown — converted with ``md_to_html_rich``.
        engagement_meta: Optional dict rendered on the cover meta card.
        show_cover: If True, render a dedicated cover page (default).
        cover_stats: Optional list of ``(label, value, meta)`` rendered as
            stat cards in the cover right column. If omitted, the function
            tries to derive 3-4 stats from ``engagement_meta``.
        engagement_id: Optional engagement id, surfaced in the cover banner.
    """
    body_html = md_to_html_rich(body_md)
    css = _stylesheet(practice_label)
    label_esc = _html_escape(practice_label)
    title_esc = _html_escape(title)
    subtitle_esc = _html_escape(subtitle)
    today = _today_pt()
    year = _now().strftime("%Y")

    if show_cover:
        meta = engagement_meta or {}
        stats = cover_stats or _derive_cover_stats(meta)
        eng_id_html = (
            f'<div class="cover-engagement-id">'
            f'Engagement {_html_escape(engagement_id)}</div>'
            if engagement_id else ""
        )
        cover_block = f"""
<section class="cover-page">
  <div class="cover-banner">
    <div class="cover-banner-top">
      <div class="cover-wordmark">ANUVIA</div>
      <div class="cover-banner-eyebrow">{label_esc}</div>
    </div>
    <h1 class="cover-title">{title_esc}</h1>
    <p class="cover-subtitle">{subtitle_esc}</p>
    {eng_id_html}
  </div>
  <div class="cover-body">
    <div>
      <p class="cover-body-eyebrow">Engagement</p>
      {_cover_meta_block(meta)}
    </div>
    <div>
      <p class="cover-body-eyebrow">Sumário em números</p>
      {stat_grid_html(stats) if stats else ''}
    </div>
  </div>
  <div class="cover-bottom">
    <div>Anuvia Cloud &amp; AI Consulting · São Paulo · {year}</div>
    <div>{today} · Documento confidencial</div>
  </div>
</section>
""".strip()
    else:
        cover_block = ""

    body_block = f"""
<section class="body-wrap">
  <p class="eyebrow body-eyebrow">ANUVIA · {label_esc}</p>
  <h1 class="body-title">{title_esc}</h1>
  <p class="body-subtitle">{subtitle_esc}</p>
  <div class="body-rule"></div>
  <div class="body-content">
{body_html}
  </div>
</section>
""".strip()

    return f"""<!DOCTYPE html>
<html lang="pt-BR">
<head>
<meta charset="utf-8">
<title>{title_esc}</title>
<link rel="preconnect" href="https://fonts.googleapis.com">
<link rel="preconnect" href="https://fonts.gstatic.com" crossorigin>
<link href="{GOOGLE_FONTS_LINK}" rel="stylesheet">
<style>{css}</style>
</head>
<body>
{cover_block}
{body_block}
</body>
</html>"""


# ---------------------------------------------------------------------------
# PPTX deck generation (light theme — matches LP, NOT proposal-sidecar dark)
# ---------------------------------------------------------------------------


# Brand tokens in pptx RGBColor form. Defined as hex tuples here and lifted
# at call time so module-level import doesn't blow up if python-pptx is
# missing (we want the module to import cleanly even without it).
_PPTX_INK_HEX = (0x1A, 0x1A, 0x1A)
_PPTX_PAPER_HEX = (0xFA, 0xFA, 0xF9)
_PPTX_WHITE_HEX = (0xFF, 0xFF, 0xFF)
_PPTX_STONE_500_HEX = (0x78, 0x71, 0x6C)
_PPTX_STONE_200_HEX = (0xE7, 0xE5, 0xE4)
_PPTX_MUTED_HEX = (0x47, 0x55, 0x69)

# 16:9 widescreen — same as proposal sidecar.
_PPTX_SLIDE_WIDTH_IN = 13.333
_PPTX_SLIDE_HEIGHT_IN = 7.5


def _import_pptx():
    """Lazy import — raises a clear error if python-pptx isn't installed."""
    from pptx import Presentation  # type: ignore
    from pptx.util import Inches, Pt, Emu  # type: ignore
    from pptx.dml.color import RGBColor  # type: ignore
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR  # type: ignore

    return {
        "Presentation": Presentation,
        "Inches": Inches,
        "Pt": Pt,
        "Emu": Emu,
        "RGBColor": RGBColor,
        "MSO_SHAPE": MSO_SHAPE,
        "PP_ALIGN": PP_ALIGN,
        "MSO_ANCHOR": MSO_ANCHOR,
    }


def _add_rect(slide, mod, left, top, width, height, color_hex):
    shp = slide.shapes.add_shape(mod["MSO_SHAPE"].RECTANGLE, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = mod["RGBColor"](*color_hex)
    shp.line.fill.background()
    shp.shadow.inherit = False
    return shp


def _add_text(
    slide,
    mod,
    text,
    left,
    top,
    width,
    height,
    *,
    size=14,
    bold=False,
    italic=False,
    color_hex=_PPTX_INK_HEX,
    font="Calibri",
    align=None,
    anchor=None,
    letter_spacing=None,
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = mod["Emu"](0)
    tf.margin_right = mod["Emu"](0)
    tf.margin_top = mod["Emu"](0)
    tf.margin_bottom = mod["Emu"](0)
    if anchor is not None:
        tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    if align is not None:
        p.alignment = align
    run = p.add_run()
    run.text = text or ""
    run.font.size = mod["Pt"](size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = mod["RGBColor"](*color_hex)
    run.font.name = font
    return tb


def _add_bullets(
    slide,
    mod,
    bullets: List[str],
    left,
    top,
    width,
    height,
    *,
    size=16,
    color_hex=_PPTX_INK_HEX,
    font="Calibri",
    bullet_char="—",
):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = mod["Emu"](0)
    tf.margin_right = mod["Emu"](0)
    tf.margin_top = mod["Emu"](0)
    tf.margin_bottom = mod["Emu"](0)
    for idx, raw in enumerate(bullets or []):
        text = (raw or "").strip()
        if not text:
            continue
        if idx == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = mod["Pt"](6)
        run = p.add_run()
        run.text = f"{bullet_char}  {text}"
        run.font.size = mod["Pt"](size)
        run.font.color.rgb = mod["RGBColor"](*color_hex)
        run.font.name = font
    return tb


def _set_speaker_notes(slide, notes: Optional[str]) -> None:
    if not notes:
        return
    try:
        slide.notes_slide.notes_text_frame.text = notes
    except Exception:  # noqa: BLE001
        # Some pptx versions are picky; non-fatal.
        pass


def _slide_chrome(slide, mod, practice_label: str, page_num: int, total: int):
    """Slim dark band at the top with ANUVIA brand left + page N/total right.

    Adds the matching bottom footer (subtle) and a thin rule below the band
    so the title clearly hangs under it. Used on every content slide.
    """
    # Slim dark band — ~10mm tall at the very top.
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](0.4),
        _PPTX_INK_HEX,
    )
    # Brand left.
    _add_text(
        slide, mod,
        f"ANUVIA · {practice_label}",
        mod["Inches"](0.5),
        mod["Inches"](0.08),
        mod["Inches"](8),
        mod["Inches"](0.3),
        size=9,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Calibri",
    )
    # Page x / y right.
    _add_text(
        slide, mod,
        f"{page_num:02d} / {total:02d}",
        mod["Inches"](11.5),
        mod["Inches"](0.08),
        mod["Inches"](1.33),
        mod["Inches"](0.3),
        size=9,
        color_hex=(0x94, 0xA3, 0xB8),
        font="Calibri",
        align=mod["PP_ALIGN"].RIGHT,
    )
    # Bottom footer left (date + brand).
    _add_text(
        slide, mod,
        f"Anuvia Cloud & AI Consulting · {_today_pt()}",
        mod["Inches"](0.5),
        mod["Inches"](7.15),
        mod["Inches"](8),
        mod["Inches"](0.3),
        size=8,
        color_hex=_PPTX_STONE_500_HEX,
    )
    # Bottom rule (thin).
    _add_rect(
        slide, mod,
        mod["Inches"](0.5),
        mod["Inches"](7.1),
        mod["Inches"](12.33),
        mod["Emu"](9525),
        _PPTX_STONE_200_HEX,
    )


def _build_cover_slide(prs, mod, *, practice_label, title, subtitle, client_name, engagement_id):
    """Cover slide — 40/60 split (dark left panel, white right panel)."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    # White right side.
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_WHITE_HEX,
    )
    # Dark left panel — 40% of slide width.
    dark_w = mod["Inches"](_PPTX_SLIDE_WIDTH_IN * 0.40)
    _add_rect(
        slide, mod,
        0, 0,
        dark_w,
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_INK_HEX,
    )
    # Wordmark top-left of dark panel.
    _add_text(
        slide, mod,
        "ANUVIA",
        mod["Inches"](0.55),
        mod["Inches"](0.55),
        mod["Inches"](4),
        mod["Inches"](0.5),
        size=20,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
    )
    # Practice eyebrow under wordmark.
    _add_text(
        slide, mod,
        practice_label.upper(),
        mod["Inches"](0.55),
        mod["Inches"](1.05),
        mod["Inches"](4.5),
        mod["Inches"](0.3),
        size=10,
        bold=True,
        color_hex=(0x94, 0xA3, 0xB8),  # STONE_400
    )
    # Thin rule under eyebrow.
    _add_rect(
        slide, mod,
        mod["Inches"](0.55),
        mod["Inches"](1.55),
        mod["Inches"](0.6),
        mod["Emu"](19050),
        _PPTX_WHITE_HEX,
    )
    # Big Playfair title — sits inside dark panel, but spills slightly into
    # the white side so the eye is led across the split.
    _add_text(
        slide, mod,
        title,
        mod["Inches"](0.55),
        mod["Inches"](2.6),
        mod["Inches"](7.5),
        mod["Inches"](2.4),
        size=40,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
    )
    if subtitle:
        _add_text(
            slide, mod,
            subtitle,
            mod["Inches"](0.55),
            mod["Inches"](5.05),
            mod["Inches"](4.5),
            mod["Inches"](0.5),
            size=13,
            color_hex=(0x94, 0xA3, 0xB8),
        )
    # Engagement id bottom of dark panel.
    if engagement_id:
        _add_text(
            slide, mod,
            f"Engagement · {engagement_id}",
            mod["Inches"](0.55),
            mod["Inches"](6.7),
            mod["Inches"](4.5),
            mod["Inches"](0.3),
            size=9,
            color_hex=(0x94, 0xA3, 0xB8),
            font="Calibri",
        )

    # Right side — client name, date, anuvia footer.
    right_x = mod["Inches"](_PPTX_SLIDE_WIDTH_IN * 0.40 + 0.5)
    right_w = mod["Inches"](_PPTX_SLIDE_WIDTH_IN * 0.60 - 1.0)

    _add_text(
        slide, mod,
        "Cliente",
        right_x,
        mod["Inches"](2.6),
        right_w,
        mod["Inches"](0.3),
        size=9,
        bold=True,
        color_hex=_PPTX_STONE_500_HEX,
    )
    _add_text(
        slide, mod,
        client_name or "—",
        right_x,
        mod["Inches"](2.9),
        right_w,
        mod["Inches"](0.6),
        size=22,
        bold=True,
        color_hex=_PPTX_INK_HEX,
        font="Playfair Display",
    )
    _add_text(
        slide, mod,
        "Entrega",
        right_x,
        mod["Inches"](3.9),
        right_w,
        mod["Inches"](0.3),
        size=9,
        bold=True,
        color_hex=_PPTX_STONE_500_HEX,
    )
    _add_text(
        slide, mod,
        f"{_today_pt()} · Auditoria FinOps · 4 semanas",
        right_x,
        mod["Inches"](4.2),
        right_w,
        mod["Inches"](0.4),
        size=13,
        color_hex=_PPTX_INK_HEX,
    )
    # Right side bottom — signature.
    _add_rect(
        slide, mod,
        right_x,
        mod["Inches"](6.4),
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN * 0.60 - 1.0),
        mod["Emu"](9525),
        _PPTX_STONE_200_HEX,
    )
    _add_text(
        slide, mod,
        "Mila Vernazza · founder@anuvia.com.br",
        right_x,
        mod["Inches"](6.55),
        right_w,
        mod["Inches"](0.3),
        size=10,
        bold=True,
        color_hex=_PPTX_INK_HEX,
    )
    _add_text(
        slide, mod,
        "Anuvia Cloud & AI Consulting · São Paulo",
        right_x,
        mod["Inches"](6.85),
        right_w,
        mod["Inches"](0.3),
        size=9,
        color_hex=_PPTX_STONE_500_HEX,
    )


def _build_section_slide(prs, mod, *, title, subtitle=None, section_num=None):
    """Full-bleed dark slate section divider. Centered Playfair title."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_INK_HEX,
    )
    # Tiny eyebrow "SEÇÃO N".
    eyebrow = (
        f"SEÇÃO {section_num:02d}" if isinstance(section_num, int) else "SEÇÃO"
    )
    _add_text(
        slide, mod,
        eyebrow,
        mod["Inches"](0.7),
        mod["Inches"](2.55),
        mod["Inches"](11.9),
        mod["Inches"](0.35),
        size=10,
        bold=True,
        color_hex=(0x94, 0xA3, 0xB8),
        align=mod["PP_ALIGN"].CENTER,
    )
    _add_text(
        slide, mod,
        title,
        mod["Inches"](0.7),
        mod["Inches"](3.1),
        mod["Inches"](11.9),
        mod["Inches"](1.6),
        size=44,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
        align=mod["PP_ALIGN"].CENTER,
    )
    if subtitle:
        _add_text(
            slide, mod,
            subtitle,
            mod["Inches"](0.7),
            mod["Inches"](4.65),
            mod["Inches"](11.9),
            mod["Inches"](0.7),
            size=16,
            color_hex=(0x94, 0xA3, 0xB8),
            align=mod["PP_ALIGN"].CENTER,
        )
    # Thin centered rule.
    _add_rect(
        slide, mod,
        mod["Inches"](6.166),
        mod["Inches"](5.5),
        mod["Inches"](1.0),
        mod["Emu"](19050),
        _PPTX_WHITE_HEX,
    )


def _build_content_slide(prs, mod, *, practice_label, title, subtitle, bullets, page_num, total):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_WHITE_HEX,
    )
    _slide_chrome(slide, mod, practice_label, page_num, total)
    # Title — sits below the dark top band.
    _add_text(
        slide, mod,
        title,
        mod["Inches"](0.7),
        mod["Inches"](0.95),
        mod["Inches"](12.0),
        mod["Inches"](0.9),
        size=28,
        bold=True,
        color_hex=_PPTX_INK_HEX,
        font="Playfair Display",
    )
    if subtitle:
        _add_text(
            slide, mod,
            subtitle,
            mod["Inches"](0.7),
            mod["Inches"](1.85),
            mod["Inches"](12.0),
            mod["Inches"](0.4),
            size=12,
            color_hex=_PPTX_MUTED_HEX,
        )
    # Bullets — indented 8mm (~0.31") from the title left, per spec.
    _add_bullets(
        slide, mod,
        bullets or [],
        mod["Inches"](1.01),  # 0.7 + 0.31 indent
        mod["Inches"](2.55),
        mod["Inches"](11.7),
        mod["Inches"](4.3),
        size=15,
        color_hex=_PPTX_INK_HEX,
    )
    return slide


def _build_two_col_slide(
    prs, mod, *,
    practice_label, title, subtitle,
    left_title, left_bullets,
    right_title, right_bullets,
    page_num, total,
):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_WHITE_HEX,
    )
    _slide_chrome(slide, mod, practice_label, page_num, total)
    _add_text(
        slide, mod,
        title,
        mod["Inches"](0.7),
        mod["Inches"](0.95),
        mod["Inches"](12.0),
        mod["Inches"](0.9),
        size=28,
        bold=True,
        color_hex=_PPTX_INK_HEX,
        font="Playfair Display",
    )
    if subtitle:
        _add_text(
            slide, mod,
            subtitle,
            mod["Inches"](0.7),
            mod["Inches"](1.85),
            mod["Inches"](12.0),
            mod["Inches"](0.4),
            size=12,
            color_hex=_PPTX_MUTED_HEX,
        )
    # Vertical rule between columns.
    _add_rect(
        slide, mod,
        mod["Inches"](6.65),
        mod["Inches"](2.4),
        mod["Emu"](9525),
        mod["Inches"](4.4),
        _PPTX_STONE_200_HEX,
    )
    # Column titles.
    if left_title:
        _add_text(
            slide, mod, left_title,
            mod["Inches"](0.7),
            mod["Inches"](2.4),
            mod["Inches"](5.7),
            mod["Inches"](0.4),
            size=10, bold=True,
            color_hex=_PPTX_STONE_500_HEX,
        )
    if right_title:
        _add_text(
            slide, mod, right_title,
            mod["Inches"](6.95),
            mod["Inches"](2.4),
            mod["Inches"](5.7),
            mod["Inches"](0.4),
            size=10, bold=True,
            color_hex=_PPTX_STONE_500_HEX,
        )
    # Bullets.
    _add_bullets(
        slide, mod,
        left_bullets or [],
        mod["Inches"](0.7),
        mod["Inches"](2.85),
        mod["Inches"](5.7),
        mod["Inches"](3.9),
        size=14,
        color_hex=_PPTX_INK_HEX,
    )
    _add_bullets(
        slide, mod,
        right_bullets or [],
        mod["Inches"](6.95),
        mod["Inches"](2.85),
        mod["Inches"](5.7),
        mod["Inches"](3.9),
        size=14,
        color_hex=_PPTX_INK_HEX,
    )
    return slide


def _build_closing_slide(prs, mod, *, practice_label, title, subtitle, bullets, page_num, total):
    """Closing slide — dark background, Mila signature + LinkedIn placeholder."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_rect(
        slide, mod,
        0, 0,
        mod["Inches"](_PPTX_SLIDE_WIDTH_IN),
        mod["Inches"](_PPTX_SLIDE_HEIGHT_IN),
        _PPTX_INK_HEX,
    )
    # Wordmark top-left.
    _add_text(
        slide, mod,
        "ANUVIA",
        mod["Inches"](0.7),
        mod["Inches"](0.55),
        mod["Inches"](4),
        mod["Inches"](0.5),
        size=18,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
    )
    _add_text(
        slide, mod,
        practice_label.upper(),
        mod["Inches"](0.7),
        mod["Inches"](1.05),
        mod["Inches"](5),
        mod["Inches"](0.3),
        size=9,
        bold=True,
        color_hex=(0x94, 0xA3, 0xB8),
    )
    # Page indicator top-right.
    _add_text(
        slide, mod,
        f"{page_num:02d} / {total:02d}",
        mod["Inches"](11.5),
        mod["Inches"](0.55),
        mod["Inches"](1.33),
        mod["Inches"](0.3),
        size=9,
        color_hex=(0x94, 0xA3, 0xB8),
        align=mod["PP_ALIGN"].RIGHT,
    )

    _add_text(
        slide, mod,
        title or "Próximos passos",
        mod["Inches"](0.7),
        mod["Inches"](2.2),
        mod["Inches"](12.0),
        mod["Inches"](1.0),
        size=36,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
    )
    if subtitle:
        _add_text(
            slide, mod, subtitle,
            mod["Inches"](0.7),
            mod["Inches"](3.3),
            mod["Inches"](12.0),
            mod["Inches"](0.5),
            size=14,
            color_hex=(0x94, 0xA3, 0xB8),
        )
    _add_bullets(
        slide, mod,
        bullets or [],
        mod["Inches"](0.7),
        mod["Inches"](4.0),
        mod["Inches"](12.0),
        mod["Inches"](2.4),
        size=14,
        color_hex=_PPTX_WHITE_HEX,
    )

    # Signature block — thin white rule then sig + LinkedIn.
    _add_rect(
        slide, mod,
        mod["Inches"](0.7),
        mod["Inches"](6.55),
        mod["Inches"](2.5),
        mod["Emu"](19050),
        _PPTX_WHITE_HEX,
    )
    _add_text(
        slide, mod,
        "Mila Vernazza",
        mod["Inches"](0.7),
        mod["Inches"](6.7),
        mod["Inches"](6),
        mod["Inches"](0.35),
        size=13,
        bold=True,
        color_hex=_PPTX_WHITE_HEX,
        font="Playfair Display",
    )
    _add_text(
        slide, mod,
        "founder · Anuvia Cloud & AI Consulting",
        mod["Inches"](0.7),
        mod["Inches"](7.05),
        mod["Inches"](6),
        mod["Inches"](0.3),
        size=9,
        color_hex=(0x94, 0xA3, 0xB8),
    )
    # Contact right column.
    _add_text(
        slide, mod,
        "mila@anuvia.com.br",
        mod["Inches"](7.0),
        mod["Inches"](6.7),
        mod["Inches"](5.8),
        mod["Inches"](0.3),
        size=10,
        color_hex=_PPTX_WHITE_HEX,
        align=mod["PP_ALIGN"].RIGHT,
    )
    _add_text(
        slide, mod,
        "linkedin.com/in/milavernazza",
        mod["Inches"](7.0),
        mod["Inches"](7.05),
        mod["Inches"](5.8),
        mod["Inches"](0.3),
        size=9,
        color_hex=(0x94, 0xA3, 0xB8),
        align=mod["PP_ALIGN"].RIGHT,
    )


async def generate_pptx_deck(
    *,
    practice_label: str,
    title: str,
    client_name: str,
    engagement_id: str,
    slides: List[dict],
) -> bytes:
    """Build a .pptx Presentation in-memory and return the binary bytes.

    ``slides`` is a list of dicts shaped like::

        {
            "type": "cover" | "section" | "content" | "two_col" | "closing",
            "title": str,
            "subtitle": str | None,
            "bullets": list[str] | None,
            "notes": str | None,
            "left_bullets": list[str] | None,
            "right_bullets": list[str] | None,
            "left_title": str | None,
            "right_title": str | None,
        }

    The first slide in the list is rendered as a cover (overriding type),
    the last as a closing slide (also overriding) so callers can pass a
    plain list of content-shaped specs and still get bookends.
    """
    mod = _import_pptx()
    prs = mod["Presentation"]()
    prs.slide_width = mod["Inches"](_PPTX_SLIDE_WIDTH_IN)
    prs.slide_height = mod["Inches"](_PPTX_SLIDE_HEIGHT_IN)

    if not slides:
        slides = [{
            "type": "content",
            "title": title,
            "subtitle": "",
            "bullets": ["(deck vazio — engenharia revisar)"],
        }]

    total = len(slides)
    section_counter = 0

    for idx, spec in enumerate(slides):
        page_num = idx + 1
        # Force bookends.
        if idx == 0:
            stype = "cover"
        elif idx == total - 1 and total > 1:
            stype = spec.get("type") or "closing"
            if stype not in ("closing", "section"):
                stype = "closing"
        else:
            stype = spec.get("type") or "content"

        s_title = spec.get("title") or ""
        s_subtitle = spec.get("subtitle") or ""
        bullets = spec.get("bullets") or []

        if stype == "cover":
            _build_cover_slide(
                prs, mod,
                practice_label=practice_label,
                title=s_title or title,
                subtitle=s_subtitle,
                client_name=client_name,
                engagement_id=engagement_id,
            )
            slide_obj = prs.slides[-1]
        elif stype == "section":
            section_counter += 1
            _build_section_slide(
                prs, mod,
                title=s_title,
                subtitle=s_subtitle,
                section_num=section_counter,
            )
            slide_obj = prs.slides[-1]
        elif stype == "two_col":
            slide_obj = _build_two_col_slide(
                prs, mod,
                practice_label=practice_label,
                title=s_title,
                subtitle=s_subtitle,
                left_title=spec.get("left_title"),
                left_bullets=spec.get("left_bullets") or [],
                right_title=spec.get("right_title"),
                right_bullets=spec.get("right_bullets") or [],
                page_num=page_num,
                total=total,
            )
        elif stype == "closing":
            _build_closing_slide(
                prs, mod,
                practice_label=practice_label,
                title=s_title or "Próximos passos",
                subtitle=s_subtitle,
                bullets=bullets,
                page_num=page_num,
                total=total,
            )
            slide_obj = prs.slides[-1]
        else:  # content
            slide_obj = _build_content_slide(
                prs, mod,
                practice_label=practice_label,
                title=s_title,
                subtitle=s_subtitle,
                bullets=bullets,
                page_num=page_num,
                total=total,
            )

        _set_speaker_notes(slide_obj, spec.get("notes"))

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue()


# ---------------------------------------------------------------------------
# Markdown deck parser — turns Claude's "### Slide N — Title" md into specs
# ---------------------------------------------------------------------------


_SLIDE_HEADER_RE = re.compile(
    r"^#{2,4}\s*(?:Slide\s*)?(\d+)?\s*[—\-:.]?\s*(.*?)\s*$",
    re.IGNORECASE,
)
_NOTES_RE = re.compile(r"^\s*\(?\s*notas?\s*:\s*(.+?)\)?\s*$", re.IGNORECASE)


def parse_deck_markdown(md: str) -> List[dict]:
    """Parse a Claude-generated deck markdown into slide specs.

    Recognised structure::

        ### Slide 1 — Capa
        - bullet 1
        - bullet 2

        (notas: speaker notes)

        ### Slide 2 — Sumário
        - bullet

    Heuristics:
      * First slide -> cover
      * Last slide -> closing
      * Slides whose title contains "Vetor N" / "Quick Win" / etc -> content
      * Notes captured from any "(notas: ...)" or "Notas: ..." line
    """
    if not md:
        return []

    slides: List[dict] = []
    current: Optional[dict] = None
    notes_lines: List[str] = []

    def _flush_notes_to_current():
        nonlocal notes_lines
        if current is not None and notes_lines:
            joined = "\n".join(line.strip() for line in notes_lines if line.strip())
            if joined:
                existing = current.get("notes") or ""
                current["notes"] = (existing + "\n\n" + joined).strip() if existing else joined
        notes_lines = []

    for raw in md.splitlines():
        line = raw.rstrip()
        stripped = line.strip()
        if not stripped:
            continue

        # New slide header?
        if stripped.startswith("#"):
            m = _SLIDE_HEADER_RE.match(stripped.lstrip("#").strip())
            title_text = stripped.lstrip("#").strip()
            # Strip leading "Slide N —" prefix from the title.
            cleaned = re.sub(
                r"^Slide\s*\d+\s*[—\-:.]?\s*", "", title_text, flags=re.IGNORECASE,
            )
            if current is not None:
                _flush_notes_to_current()
                slides.append(current)
            current = {
                "type": "content",
                "title": cleaned or title_text,
                "subtitle": "",
                "bullets": [],
                "notes": "",
            }
            continue

        # Speaker notes line — "(notas: ...)" or "Notas: ..."
        n = _NOTES_RE.match(stripped)
        if n:
            notes_lines.append(n.group(1).strip())
            continue

        # Bullet?
        bm = re.match(r"^[-*•]\s+(.*)$", stripped)
        if bm and current is not None:
            bullet = bm.group(1).strip()
            # Strip markdown bold markers — they don't render in PPTX runs.
            bullet = re.sub(r"\*\*(.+?)\*\*", r"\1", bullet)
            bullet = re.sub(r"`([^`]+)`", r"\1", bullet)
            current.setdefault("bullets", []).append(bullet)
            continue

        # Free-floating text → subtitle (first line) or extra notes.
        if current is not None:
            if not current.get("subtitle"):
                # Strip md formatting.
                txt = re.sub(r"\*\*(.+?)\*\*", r"\1", stripped)
                txt = re.sub(r"`([^`]+)`", r"\1", txt)
                # Don't promote bullets-disguised-as-text to subtitle.
                current["subtitle"] = txt[:160]
            else:
                notes_lines.append(stripped)

    if current is not None:
        _flush_notes_to_current()
        slides.append(current)

    # Assign types: first=cover, last=closing.
    total = len(slides)
    for idx, s in enumerate(slides):
        title_low = (s.get("title") or "").lower()
        if idx == 0:
            s["type"] = "cover"
        elif idx == total - 1 and total > 1:
            # Closing if the title hints at it, otherwise still closing
            # (consistent bookends).
            s["type"] = "closing"
        else:
            # Section dividers when title is purely a thematic break.
            if any(k in title_low for k in ("sumário", "agenda", "índice")):
                s["type"] = "content"
            else:
                s["type"] = "content"

    return slides
